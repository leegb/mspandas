import pptx
import pandas as pd
import numpy as np

class Handler():
    """Convenience tools for accessing PPT templates using python-pptx.

    Methods
    -------
    map_layouts: Create dictionary object of template layouts in slide master from ppt object, where keys are layout names.
    map_shapes: Create dictionary object of slide shapes in template layout from layout object, where keys are shape names.
    """

    def map_layouts(self, ppt, verbose=False):
        """Create dictionary object of template layouts in slide master from ppt object, where keys are layout names.

        Parameters
        ----------
        ppt: ppt.Presentation
            Powerpoint presentation object.

        Returns
        -------
        layout_map: dict
            Dictionary of ppt layout objects where keys are layout names from slide master.
        """
        layout_map = {}
        for slide in ppt.slide_layouts:
            layout_map[slide.name] = slide
            if verbose:
                print(slide.name)
        return layout_map

    def map_shapes(self, layout, verbose=False):
        """Create dictionary object of slide shapes in template layout from layout object, where keys are shape names.

        Parameters
        ----------
        layout: ppt.slide.SlideLayout
            Slide layout object.

        Returns
        -------
        shape_map: dict
            Dictionary of slide shape objects where keys are shape names from template layout.
        """
        shape_map = {}
        for shape in layout.shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                shape_map[shape.name] = phf.idx
                if verbose:
                    print('{} index: {}, type: {}'.format(shape.name, phf.idx, phf.type))
        return shape_map


class Table():
    """
    Abstract table class to convert a Pandas DataFrame into a PowerPoint table shape.

    Attributes
    ----------
    shape: pptx.shapes.base.BaseShape
        pptx shape object, placeholder for table graphicframe.
    df: pd.DataFrame
        Pandas DataFrame to
         be converted into table.

    Parameters
    ----------
    header: bool
        Whether or not to include DataFrame's header in table representation.
    index: bool
        Whether or not to include DataFrame's index in table representation.
    keep_names: str
        Priority for keeping names on axis when transforming. If table includes header and index, and DataFrame names specified for both axis, will keep names based on priority.
        Specificy 'index' for column-wise totals, 1 or 'columns' for row-wise totals.
    na_rep: int,float,str
        String value for representing null values.
    dtype_format: dict
        Map of numpy data types to string format.
    font_size: int
        Cell text font size. For more, see pptx.util.Pt
    font_color: tuple
        Cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
    font_name:
        Cell text font name, for example 'Arial'.
    column_totals: bool
        Whether or not to include column totals in table representation.
    column_totals_label: str
        Index label for totals row, can be string or tuple for multiindex.
    column_totals_method: str/func
        Function to use for aggregating the columns.
    column_totals_aggmap: dict
        Map of column names and aggregation method to be applied when totaling columns. Example: {'a':'mean', 'b':'mean'}. Default will sum all data.
	row_totals: bool
        Whether or not to include row totals in table representation.
    row_totals_label: str
        Column label for totals column, can be string or tuple for multiindex header.
    row_totals_method: str/func
        Function to use for aggregating the rows.
    row_totals_aggmap: dict
        Map of index names and aggregation method to be applied when totaling columns. Example: {0:'mean', 1:'mean'}. Default will sum all data.
    fill_header: bool
        Whether or not to fill the cell backgound color of table header rows.
    bold_header: bool
        Whether or not to bold the text in all table header rows.
    header_font_size: int
        Header cell text font size. For more, see pptx.util.Pt
    header_font_color: tuple
        Header cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
    fill_index: bool
        Whether or not to fill the cell backgound color of table index columns.
    bold_index: bool
        Whether or not to bold the text in all table index columns.
    index_font_size: int
        Index cell text font size. For more, see pptx.util.Pt
    index_font_color: tuple
        Index cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
    fill_color: tuple or pptx.enum.dml.MSO_THEME_COLOR
        Color to fill cell background. Must be RGB code as tuple of 3 integers, or instance of pptx.enum.dml.MSO_THEME_COLOR.
    row_banding: bool
        Whether or not turn on default PowerPoint horizontal banding styles.
    column_banding: bool
        Whether or not turn on default PowerPoint vertical banding styles.
    first_row: bool
        Whether of not to turn on default PowerPoint styles for first row.
    first_col: bool
        Whether of not to turn on default PowerPoint styles for first col.
    last_row: bool
        Whether of not to turn on default PowerPoint styles for last row.
    last_row: bool
        Whether of not to turn on default PowerPoint styles for last row.
    merge_indices: bool
        Whether or not to merge table indices (rows in index columns or columns in header rows) by combining adjacent equal values.
    center_merge: str
        Whether or not to center the paragraph text after alignment.

    Methods
    -------
    add_totals: Aggregate data in DataFrame, applied as column-wise or row-wise by axis argument.
    transform: Transofrm DataFrame to mirror output representation.
    format_index: Format DataFrame index values as unicode strings.
    format_values: Format DataFrame values as unicode strings.
    insert_table: Insert a graphic frame table object into the table shape.
    style_index: Apply styles to DataFrame index in PowerPoint table in place.
    style_table: Apply styles to PowerPoint table in place.
    convert: Perform the conversion from DataFrame values to table cells.
    """

    def __init__(self, shape, df, **kwargs):
        self.shape = shape
        self.df = df

        self.header = kwargs.pop('header',False)
        self.index = kwargs.pop('index',False)
        self.keep_names = kwargs.pop('keep_names','columns')

        self.na_rep = kwargs.pop('na_rep',' ')
        self.dtype_format = kwargs.pop('dtype_format',None)

        self.font_size = kwargs.pop('font_size',None)
        self.font_color = kwargs.pop('font_color',None)
        self.font_name = kwargs.pop('font_name',None)

        self.column_totals = kwargs.pop('column_totals',False)
        self.column_totals_label = kwargs.pop('column_totals_label','Total')
        self.column_totals_method = kwargs.pop('column_totals_method',np.sum)
        self.column_totals_aggmap = kwargs.pop('column_totals_aggmap',{})
        self.row_totals = kwargs.pop('row_totals',False)
        self.row_totals_label = kwargs.pop('row_totals_label','Total')
        self.row_totals_method = kwargs.pop('row_totals_method',np.sum)
        self.row_totals_aggmap = kwargs.pop('row_totals_aggmap',{})

        self.fill_header = kwargs.pop('fill_header',True)
        self.bold_header = kwargs.pop('bold_header',True)
        self.header_font_size = kwargs.pop('header_font_size',None)
        self.header_font_color = kwargs.pop('header_font_color',None)
        self.fill_index = kwargs.pop('fill_index',False)
        self.bold_index = kwargs.pop('bold_index',True)
        self.index_font_size = kwargs.pop('index_font_size',None)
        self.index_font_color = kwargs.pop('index_font_color',None)
        self.fill_color = kwargs.pop('fill_color',pptx.enum.dml.MSO_THEME_COLOR.ACCENT_1)

        self.row_banding = kwargs.pop('row_banding',True)
        self.column_banding = kwargs.pop('column_banding',False)
        self.first_row = kwargs.pop('first_row',False)
        self.first_col = kwargs.pop('first_col',False)
        self.last_row = kwargs.pop('last_row',False)
        self.last_col = kwargs.pop('last_col',False)

        self.merge_indices = kwargs.pop('merge_indices',True)
        self.center_merge = kwargs.pop('center_merge',True)

    def add_totals(self, **kwargs):
        """Aggregate data in DataFrame, applied as column-wise or row-wise by axis argument.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to have columns totaled.
        totals_label: str/tuple
            Index (or columns) label for totals row (or column), can be string or tuple for multiindex.
        totals_method: str/func
            Function to use for aggregating data. Default will apply np.sum.
        totals_aggmap: dict
            Map of column (or index) names and aggregation method to be applied when totaling. Example: {'a':'mean', 'b':'mean'}. Takes priority over total_method.
        axis: int
            Orientation, based on axis of DataFrame index, for totaling. 0 or 'index' for column-wise totals, 1 or 'columns' for row-wise totals.

        Returns
        -------
        data: pandas.DataFrame
            DataFrame with new row as column totals.
        """
        data = kwargs.pop('data',self.df.copy())
        totals_label = kwargs.pop('totals_label','Totals')
        totals_method = kwargs.pop('totals_method',np.sum)
        totals_aggmap = kwargs.pop('totals_aggmap',{})
        axis = kwargs.pop('axis',0)
        if axis in [1,'columns']:
            data = data.T
        elif axis in [0,'index']:
            pass
        else:
            raise(ValueError('Incorrect value for axis. Use 0 or "index" for column-wise totals, 1 or "columns" for row-wise totals.'))
        names = list(data.index.names)
        ordered_columns = list(data.columns)
        # totals
        c_totals = data.fillna(0).agg({col:totals_aggmap.get(col,totals_method) for col in data.columns}).rename(totals_label)
        c_totals = c_totals.to_frame().T
        # create multiindex if needed
        for i in range(data.index.nlevels-1):
            c_totals['dummy'] = ' '
            c_totals = c_totals.set_index('dummy',append=True)
            c_totals.index.names = [None]*len(c_totals.index.names)
        try:
        	data = pd.concat([data, c_totals], axis=0)
        except TypeError:
        	# df index is categorical
        	# add label as category and append
        	data.index = data.index.add_categories(totals_label)
        	data = pd.concat([data, c_totals], axis=0)
        data = data.reindex_axis(ordered_columns, axis=1)
        data.index.names = names
        if axis in [1,'columns']:
            return data.T
        else:
            return data

    def transform(self, **kwargs):
        """Transofrm DataFrame into output representation with presentation options.

        Optional header, index or totals are converted into DataFrame row values.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to be formatted.
        header: bool
            Whether or not to include DataFrame's header in table representation.
        index: bool
            Whether or not to include DataFrame's index in table representation.
        keep_names: str
            Priority for keeping names on axis when transforming. If table includes header and index, and DataFrame names specified for both axis, will keep names based on priority.
            Specificy 'index' for column-wise totals, 1 or 'columns' for row-wise totals.

        Returns
        -------
        data: pandas.DataFrame
            Transformed DataFrame including optional header, index and totals.

        Notes:
        ------
        When header is True and DataFrame column names are strings, \
        Series of non-object dtypes will be converted to having dtypes of object.
        """
        data = kwargs.pop('data',self.df.copy())
        header = kwargs.pop('header',self.header)
        index = kwargs.pop('index',self.index)
        keep_names = kwargs.pop('keep_names',self.keep_names)
        if keep_names == 'index':
            if index:
                data = data.reset_index()
            if header:
                data = data.T.reset_index().T
        elif keep_names == 'columns':
            if header:
                data = data.T.reset_index().T
            if index:
                data = data.reset_index()
        else:
            raise(ValueError('Incorrect value for keep_names. Use "index", or "columns"'))
        return data

    def format_index(self, **kwargs):
        """Format DataFrame index values as unicode strings.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to be formatted.
        dtype_format: dict
            Map of numpy data types to string format.
        axis: int
            Axis of DataFrame index to be formatted. 0 or 'index' for index, 1 or 'columns' for header.

        Returns
        -------
        data: pandas.DataFrame
            Formatted DataFrame where all index values are dtype np.unicode.
        """
        data = kwargs.pop('data',self.df.copy())
        dtype_format = kwargs.pop('dtype_format',self.dtype_format)
        axis = kwargs.pop('axis',0)
        if axis in [1,'columns']:
            index = data.columns
        elif axis in [0,'index']:
            index = data.index
        else:
            raise(ValueError('Incorrect value for axis. Use 0 or "index" for index, 1 or "columns" for header.'))
        names=index.names
        index_vals=[]
        for n in range(index.nlevels):
            vals = index.get_level_values(n)
            if dtype_format is not None:
                for dtype,fmt in dtype_format.items():
                    if np.issubdtype(vals.dtype,dtype):
                        if np.issubdtype(dtype,np.datetime64):
                            vals = vals.strftime(fmt)
                        elif np.issubdtype(dtype,np.number):
                            vals = vals.format(fmt)
                        else:
                            raise NotImplementedError('Not able to convert values of dtype {} to strings.\
                            Convert manually in your DataFrame before passing into Table()'.format(dtype))
            vals = vals.astype(np.unicode)
            index_vals.append(vals)
        index = pd.MultiIndex.from_arrays(index_vals,names=names)
        if axis in [1,'columns']:
            data.columns = index
        elif axis in [0,'index']:
            data.index = index
        return data

    def format_values(self, **kwargs):
        """Format DataFrame values as unicode strings.

        Parameters
        ----------
        data: pandas.DataFrame
            DataFrame to be formatted.
        dtype_format: dict
            Map of numpy data types to string format.

        Returns
        -------
        data: pandas.DataFrame
            Formatted DataFrame where all values are dtype np.unicode.
        """
        data = kwargs.pop('data',self.df.copy())
        dtype_format = kwargs.pop('dtype_format',self.dtype_format)
        if dtype_format is not None:
            for dtype,fmt in dtype_format.items():
                for col,x in data.iteritems():
                    if np.issubdtype(x.dtype,dtype):
                        if np.issubdtype(dtype,np.datetime64):
                            data.loc[:,col] = x.dt.strftime(fmt)
                        elif np.issubdtype(dtype,np.number):
                            data.loc[:,col] = x.apply(fmt.format)
                        else:
                            raise NotImplementedError('Not able to convert values of dtype {} to strings.\
                            Convert manually in your DataFrame before passing into Table()'.format(dtype))
        data = data.fillna(self.na_rep)
        data = data.astype(np.unicode)
        return data

    def insert_table(self, **kwargs):
        """Insert a graphic frame table object into the table shape.

        Parameters
        ----------
        overwrite: bool
            Force inserting of a table into shape. Overwrites any existing table.

        """
        overwrite = kwargs.pop('overwrite',False)
        rows,cols = self.transform().shape
        if self.column_totals:
            rows+=1
        if self.row_totals:
            cols+=1
        if not self.shape.has_table or overwrite and not self.shape.shape_id == None:
            self.shape = self.shape.insert_table(rows=rows,
                                                 cols=cols)
        else:
            raise Warning('Shape object already contains a table graphic frame.')

    def style_index(self, **kwargs):
        """Apply styles to DataFrame index in PowerPoint table in place.

        Parameters
        ----------
        fill: bool
            Whether or not to fill the cell backgound color.
        bold: bool
            Whether or not to bold the text.
        font_size: int
            Cell text font size. For more, see pptx.util.Pt
        font_color: tuple
            Cell text font color. Must be RGB code as tuple of 3 integers, or HEX code as string. For more see pptx.dml.color.RGBColor
        fill_color: tuple or pptx.enum.dml.MSO_THEME_COLOR
            Color to fill cell background. Must be RGB code as tuple of 3 integers, or instance of pptx.enum.dml.MSO_THEME_COLOR.
        merge_indices: bool
			Whether or not to merge table indices (rows in index columns or columns in header rows) by combining adjacent equal values.
        center_merge: str
            Whether or not to center the paragraph text after merge.
        axis: int
            Axis of DataFrame index to be formatted. 0 or 'index' for index, 1 or 'columns' for header.

        """
        table = self.shape.table
        fill = kwargs.pop('fill',True)
        bold = kwargs.pop('bold',False)
        font_size = kwargs.pop('font_size',None)
        font_color = kwargs.pop('font_color',None)
        fill_color = kwargs.pop('fill_color',self.fill_color)
        merge_indices = kwargs.pop('merge_indices',self.merge_indices)
        center_merge = kwargs.pop('center_merge',self.center_merge)
        axis = kwargs.pop('axis',0)
        data = self.df.copy()
        rows,cols = self.transform().shape
        if axis in [1,'columns']:
            index = data.columns
            axis = 1
            offset = cols - len(index)
            numcells = cols+1 if self.row_totals else cols
        elif axis in [0,'index']:
            index = data.index
            axis = 0
            offset = rows - len(index)
            numcells = rows+1 if self.column_totals else rows
        else:
            raise(ValueError('Incorrect value for axis. Use 0 or "index" for index, 1 or "columns" for header.'))
        for n in range(index.nlevels):
            merge = False
            for i in range(numcells):
                if axis==0:
                    c = table.cell(i,n)
                else:
                    c = table.cell(n,i)
                ### Start Section on Auto-Merging (unstable feature)
                if merge_indices and i >= offset:
                    # table loc to dataframe loc
                    j = i-offset
                    if j < len(index)-1:
                        equal = True if index.get_level_values(n)[j] == index.get_level_values(n)[j+1] else False
                        if not merge and equal:
                            merge = True
                            mergestart = i
                            # default to maximum merge span
                            mergespan = len(index)-j
                            # calc actual mergespan using next non-equal value
                            if len(index) > j+1:
                                for k,d in enumerate(index.get_level_values(n)[j+1:]):
                                    if not index.get_level_values(n)[j] == d:
                                        mergespan = k+1
                                        break
                            if axis in [0,'rows']:
                                c._tc.set('rowSpan', str(mergespan))
                            else:
                                c._tc.set('gridSpan', str(mergespan))
                            if center_merge:
                                tf = c.text_frame
                                p = tf.paragraphs[0]
                                p.alignment = pptx.enum.text.PP_ALIGN.__dict__['CENTER']
                        elif merge and not equal:
                            # stop merge
                            if axis in [0,'rows']:
                                c._tc.set('vMerge', '1')
                            else:
                                c._tc.set('hMerge', '1')
                            merge = False
                    elif merge and i > mergestart and i < mergestart+mergespan-1:
                        # continue merge
                        if axis in [0,'rows']:
                            c._tc.set('vMerge', '1')
                        else:
                            c._tc.set('hMerge', '1')
                    elif merge and i > mergestart and i == mergestart+mergespan-1:
                        # end of merge
                        if axis in [0,'rows']:
                            c._tc.set('vMerge', '1')
                        else:
                            c._tc.set('hMerge', '1')
                        merge = False
                ### End Section on Auto-Merging
                if fill:
                    c.fill.solid()
                    if isinstance(fill_color,pptx.enum.base.EnumValue):
                        c.fill.fore_color.theme_color = fill_color
                    elif isinstance(fill_color,tuple):
                        c.fill.fore_color.rgb = pptx.dml.color.RGBColor(*fill_color)
                    elif isinstance(fill_color,str):
                        c.fill.fore_color.rgb = pptx.dml.color.RGBColor.from_string(fill_color) if not fill_color.startswith('#') else pptx.dml.color.RGBColor.from_string(fill_color[1:])
                    else:
                        raise ValueError('Incorrect value for fill_color. \
                        Please provide one of RGB code as `tuple` of 3 integers, HEX code as string, or an instance of `pptx.enum.dml.MSO_THEME_COLOR`')
                # apply text formatting for each cell
                tf = c.text_frame
                p = tf.paragraphs[0]
                r = p.runs[0]
                if not font_size == None:
                    r.font.size = pptx.util.Pt(font_size)
                if not font_color == None:
                    if isinstance(font_color,tuple):
                        r.font.color.rgb = pptx.dml.color.RGBColor(*font_color)
                    elif isinstance(font_color,str):
                        r.font.color.rgb = pptx.dml.color.RGBColor.from_string(font_color[1:]) if font_color.startswith('#') else pptx.dml.color.RGBColor.from_string(font_color)
                    else:
                        raise ValueError('Incorrect value for font_color. \
                        Please provide one of RGB code as `tuple` of 3 integers, or a HEX code as string')
                if bold:
                    r.font.bold = True

    def style_table(self, **kwargs):
        """Apply styles to PowerPoint table in place.

        Parameters
        ----------
        header: bool
            Whether or not to include DataFrame's header in table representation.
        fill_header: bool
            Whether or not to fill the cell backgound color of table header rows.
        bold_header: bool
            Whether or not to bold the text in all table header rows.
        index: bool
            Whether or not to include DataFrame's index in table representation.
        fill_index: bool
            Whether or not to fill the cell backgound color of table index columns.
        bold_index: bool
            Whether or not to bold the text in all table index columns.
        fill_color: tuple or pptx.enum.dml.MSO_THEME_COLOR
            Color to fill cell background. Must be RGB code as tuple of 3 integers, or instance of pptx.enum.dml.MSO_THEME_COLOR.
        row_banding: bool
            Whether or not turn on default PowerPoint horizontal banding styles.
        column_banding: bool
            Whether or not turn on default PowerPoint vertical banding styles.
        first_row: bool
            Whether of not to turn on default PowerPoint styles for first row.
        first_col: bool
            Whether of not to turn on default PowerPoint styles for first col.
        last_row: bool
            Whether of not to turn on default PowerPoint styles for last row.
        last_row: bool
            Whether of not to turn on default PowerPoint styles for last row.
        merge_indices: bool
			Whether or not to merge table indices (rows in index columns or columns in header rows) by combining adjacent equal values.
        center_merge: str
            Whether or not to center the paragraph text after merge.

        Notes
        -----
        We do not apply the PowerPoint styles, controlling overarching table theme, as they are not currently supported by python-pptx (see https://github.com/scanny/python-pptx/issues/27)
        Instead we apply the logical DataFrame styling, emphasizing header, index and totals with bold text or filled backgound
        """
        table =  self.shape.table
        fill_header = kwargs.pop('fill_header',self.fill_header)
        bold_header = kwargs.pop('bold_header',self.bold_header)
        header_font_size = kwargs.pop('header_font_size',self.header_font_size)
        header_font_color = kwargs.pop('header_font_size',self.header_font_color)
        fill_index = kwargs.pop('fill_index',self.fill_index)
        bold_index = kwargs.pop('bold_index',self.bold_index)
        index_font_size = kwargs.pop('index_font_size',self.index_font_size)
        index_font_color = kwargs.pop('index_font_color',self.index_font_color)
        fill_color = kwargs.pop('fill_color',self.fill_color)
        row_banding = kwargs.pop('row_banding',self.row_banding)
        column_banding = kwargs.pop('column_banding',self.column_banding)
        first_row = kwargs.pop('first_row',self.first_row)
        first_col = kwargs.pop('first_col',self.first_col)
        last_row = kwargs.pop('last_row',self.last_row)
        last_col = kwargs.pop('last_col',self.last_col)
        merge_indices = kwargs.pop('merge_indices',self.merge_indices)
        center_merge = kwargs.pop('center_merge',self.center_merge)
        if (fill_header or bold_header or not header_font_size == None or not header_font_color == None) and not self.header:
            raise ValueError("Cannot style DataFrame header when table.header attribute is False. \
            First set `table.header = True` to include the DataFrame header in the output table.")
        elif self.header:
            self.style_index(axis=1,
                             fill=fill_header,
                             bold=bold_header,
                             font_size=header_font_size,
                             font_color=header_font_color,
                             fill_color=fill_color,
                             merge_indices=merge_indices,
                             center_merge=center_merge)
        if (fill_index or bold_index or not index_font_size == None or not index_font_color == None) and not self.index:
            raise ValueError("Cannot style DataFrame index when table.index attribute is False. \
            First set `table.header = True` to include the DataFrame header in the output table.")
        elif self.index:
            self.style_index(axis=0,
                             fill=fill_index,
                             bold=bold_index,
                             font_size=index_font_size,
                             font_color=index_font_color,
                             fill_color=fill_color,
                             merge_indices=merge_indices,
                             center_merge=center_merge)
        table.horz_banding = row_banding
        table.vert_banding = column_banding
        table.first_row = first_row
        table.first_col = first_col
        table.last_row = last_row
        table.last_col = last_col

    def convert(self, **kwargs):
        """Perform the conversion from DataFrame values to table cells.

        Returns
        -------
        data: pandas.DataFrame
            Formatted and transformed DataFrame exported to PowerPoint.

        Notes
        -----
        This method calls the individual processing methods in sequence, \
        then builds the PowerPoint table by inserting values cell by cell.
        """
        font_size = self.font_size
        font_color = self.font_color
        font_name = self.font_name
        self.insert_table()
        table = self.shape.table
        data = self.df.copy()
        data = self.format_index(data=data, axis=0)
        data = self.format_index(data=data, axis=1)
        if self.column_totals:
            data = self.add_totals(data=data, axis=0,
            totals_label=self.column_totals_label,
            totals_method=self.column_totals_method,
            totals_aggmap=self.column_totals_aggmap)
        if self.row_totals:
            data = self.add_totals(data=data, axis=1,
            totals_label=self.row_totals_label,
            totals_method=self.row_totals_method,
            totals_aggmap=self.row_totals_aggmap)
        data = self.format_values(data=data)
        data = self.transform(data=data)
        data = data.fillna(self.na_rep)
        # add data to table cells one value at a time
        for (row,col),x in np.ndenumerate(data.values):
            c = table.cell(row,col)
            c.text = x if isinstance(x,str) else str(x)
            # apply table text formatting for each cell
            tf = c.text_frame
            try:
                p = tf.paragraphs[0]
            except IndexError:
                p = tf.add_paragraph()
            try:
                r = p.runs[0]
            except IndexError:
                r = p.add_run()
            if not font_size == None:
                r.font.size = pptx.util.Pt(font_size)
            if not font_color == None:
                if isinstance(font_color,tuple):
                    r.font.color.rgb = pptx.dml.color.RGBColor(*font_color)
                elif isinstance(font_color,str):
                    r.font.color.rgb = pptx.dml.color.RGBColor.from_string(font_color[1:]) if font_color.startswith('#') else pptx.dml.color.RGBColor.from_string(font_color)
                else:
                    raise ValueError('Incorrect value for font_color. \
                    Please provide one of RGB code as `tuple` of 3 integers, or a HEX code as string')
            if not font_name == None:
                r.font.name = font_name
        self.style_table()
        return data
